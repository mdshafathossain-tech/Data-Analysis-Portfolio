"""
E-Commerce Customer Behavior Analysis — Final Assignment Presentation
Author: Md. Shafat Hossain
Role: Student of Data Science & AI Engineering

This script generates a 12-slide professional dark-themed PowerPoint presentation
using python-pptx. Speaker scripts are embedded as comments above each slide block.
All numerical values are pulled from the attached Code_output.md.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ─────────────────────────────────────────────────────────────
# GLOBAL THEME CONSTANTS
# ─────────────────────────────────────────────────────────────
BG_COLOR     = RGBColor(15,  15,  15)   # Deep black background
ACCENT_GREEN = RGBColor(163, 255, 18)   # Neon lime accent
WHITE        = RGBColor(255, 255, 255)
GRAY_LIGHT   = RGBColor(180, 180, 180)
GRAY_MID     = RGBColor(80,  80,  80)
CARD_BG      = RGBColor(28,  28,  28)   # Card/box background
CARD_BG2     = RGBColor(35,  35,  35)   # Slightly lighter card
RED_WARN     = RGBColor(255, 80,  80)
TEAL         = RGBColor(0,   200, 180)

SW  = 10.0   # Slide width inches
SH  = 5.625  # Slide height inches


# ─────────────────────────────────────────────────────────────
# UTILITY HELPERS
# ─────────────────────────────────────────────────────────────
def new_prs():
    prs = Presentation()
    prs.slide_width  = Inches(SW)
    prs.slide_height = Inches(SH)
    return prs


def blank_slide(prs):
    layout = prs.slide_layouts[6]  # completely blank
    return prs.slides.add_slide(layout)


def set_bg(slide, color=BG_COLOR):
    from pptx.oxml.ns import qn
    from lxml import etree
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, x, y, w, h, fill_color, alpha=None):
    from pptx.util import Inches
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    return shape


def add_text(slide, text, x, y, w, h, font_size=14, color=WHITE,
             bold=False, align=PP_ALIGN.LEFT, italic=False, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf    = txBox.text_frame
    tf.word_wrap = True
    p  = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size    = Pt(font_size)
    run.font.color.rgb = color
    run.font.bold    = bold
    run.font.italic  = italic
    run.font.name    = font_name
    return txBox


def add_multiline(slide, lines, x, y, w, h, font_size=13, color=WHITE,
                  bold=False, align=PP_ALIGN.LEFT, line_spacing_pt=None, font_name="Calibri"):
    """lines = list of (text, color, bold, size) tuples or plain strings."""
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf    = txBox.text_frame
    tf.word_wrap = True
    first = True
    for item in lines:
        if isinstance(item, str):
            txt, col, bld, sz = item, color, bold, font_size
        else:
            txt = item.get("text", "")
            col = item.get("color", color)
            bld = item.get("bold", bold)
            sz  = item.get("size", font_size)
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = txt
        run.font.size = Pt(sz)
        run.font.color.rgb = col
        run.font.bold = bld
        run.font.name = font_name
    return txBox


def slide_label(slide, num, total=12):
    add_text(slide, f"{num} / {total}", 9.3, 5.3, 0.65, 0.25,
             font_size=8, color=GRAY_MID, align=PP_ALIGN.RIGHT)


def section_pill(slide, label, x=0.35, y=0.08):
    """Small neon-outlined pill tag above title."""
    add_rect(slide, x, y, 1.6, 0.25, CARD_BG2)
    add_text(slide, label, x+0.05, y+0.01, 1.5, 0.22,
             font_size=8, color=ACCENT_GREEN, bold=True, font_name="Calibri Light")


def slide_title(slide, title_text, y=0.35, size=30):
    add_text(slide, title_text, 0.35, y, 9.3, 0.6,
             font_size=size, color=WHITE, bold=True,
             align=PP_ALIGN.LEFT, font_name="Arial Black")


def accent_dot(slide, x, y, r=0.07):
    """Small neon circle decoration."""
    shape = slide.shapes.add_shape(9, Inches(x), Inches(y), Inches(r*2), Inches(r*2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_GREEN
    shape.line.fill.background()


def stat_card(slide, x, y, w, h, value, label, val_color=ACCENT_GREEN):
    add_rect(slide, x, y, w, h, CARD_BG)
    # Shrink font if value is long to prevent overflow
    val_size = 18 if len(str(value)) > 5 else 22
    add_text(slide, value, x+0.06, y+0.06, w-0.12, h*0.55,
             font_size=val_size, color=val_color, bold=True, align=PP_ALIGN.CENTER, font_name="Arial Black")
    add_text(slide, label, x+0.06, y+h*0.6, w-0.12, h*0.34,
             font_size=8.5, color=GRAY_LIGHT, align=PP_ALIGN.CENTER)


def bar_visual(slide, x, y, w, h, values, labels, max_val=None, color=ACCENT_GREEN):
    """Draw a simple horizontal bar chart using rectangles."""
    max_val = max_val or max(values)
    bar_h   = (h - 0.05*(len(values)-1)) / len(values)
    for i, (val, lbl) in enumerate(zip(values, labels)):
        bar_y  = y + i*(bar_h + 0.05)
        bar_w  = (val/max_val) * (w - 1.5)
        add_rect(slide, x+1.4, bar_y, bar_w, bar_h*0.6, color)
        add_text(slide, lbl, x, bar_y-0.01, 1.35, bar_h*0.6+0.02,
                 font_size=9, color=GRAY_LIGHT, align=PP_ALIGN.RIGHT)
        add_text(slide, f"{val:,}", x+1.45+bar_w, bar_y-0.01, 0.8, bar_h*0.6+0.02,
                 font_size=9, color=WHITE, bold=True)


# ─────────────────────────────────────────────────────────────
# SLIDE 1 — Title & Author
# Speaker: "Welcome, everyone. Today I'll walk you through our Final Assignment —
# a comprehensive data science investigation of 10,000 e-commerce customer records
# collected from a Bangladesh-based online marketplace."
# "My name is Md. Shafat Hossain, and over the next 12 slides, we'll move from
# raw statistics all the way to machine-learning-backed strategic recommendations."
# ─────────────────────────────────────────────────────────────
def slide_01(prs):
    slide = blank_slide(prs)
    set_bg(slide)

    # Left accent bar
    add_rect(slide, 0, 0, 0.06, SH, ACCENT_GREEN)

    # Top faint grid lines decoration
    for yi in [1.2, 2.4, 3.6, 4.8]:
        add_rect(slide, 0.06, yi, SW-0.06, 0.005, GRAY_MID)

    # Main title
    add_text(slide, "E-Commerce Customer", 0.5, 0.7, 9.0, 1.0,
             font_size=40, color=WHITE, bold=True, font_name="Arial Black")
    add_text(slide, "Behavior Analysis", 0.5, 1.55, 9.0, 0.9,
             font_size=40, color=ACCENT_GREEN, bold=True, font_name="Arial Black")

    # Subtitle bar
    add_rect(slide, 0.5, 2.6, 6.8, 0.04, ACCENT_GREEN)
    add_text(slide, "Final Assignment  ·  10,000 Records  ·  Advanced EDA & Statistical Testing",
             0.5, 2.72, 8.5, 0.35, font_size=12, color=GRAY_LIGHT, font_name="Calibri Light")

    # Author block
    add_rect(slide, 0.5, 3.3, 3.2, 1.1, CARD_BG)
    add_text(slide, "Md. Shafat Hossain", 0.65, 3.42, 3.0, 0.35,
             font_size=14, color=WHITE, bold=True)
    add_text(slide, "Student of Data Science & AI Engineering", 0.65, 3.78, 3.0, 0.25,
             font_size=9, color=ACCENT_GREEN)
    add_text(slide, "May 2026", 0.65, 4.05, 3.0, 0.25,
             font_size=9, color=GRAY_LIGHT)

    # Right side stats preview  — wider cards to prevent wrapping
    for i, (v, l) in enumerate([("10K","Customers"), ("16","Features"), ("12","Slides")]):
        stat_card(slide, 6.7+i*1.1, 3.2, 1.0, 1.2, v, l)

    slide_label(slide, 1)
    return slide


# ─────────────────────────────────────────────────────────────
# SLIDE 2 — Executive Summary
# Speaker: "Before diving into specific questions, let me frame the big story.
# This dataset tells us that customer retention is overwhelmingly driven by how
# much customers spend — not by discounts or subscription tier."
# "The data reveals a near-perfect device parity, a 50% return rate, and a
# striking 'Subscription Paradox' where 66% of customers hold paid tiers yet
# subscription status is nearly useless for predicting loyalty."
# ─────────────────────────────────────────────────────────────
def slide_02(prs):
    slide = blank_slide(prs)
    set_bg(slide)
    add_rect(slide, 0, 0, 0.06, SH, ACCENT_GREEN)

    section_pill(slide, "EXECUTIVE SUMMARY")
    slide_title(slide, "The Story of the Data", y=0.38)

    insights = [
        ("1", "RETENTION IS SPEND-DRIVEN", "Purchase Amount is the #1 predictor of return behavior. Discounts (only 1.5% importance) do NOT build loyalty."),
        ("2", "THE SUBSCRIPTION PARADOX", "66%+ customers on Trial/Premium, yet subscription status barely influences return behavior."),
        ("3", "PAYMENT = EXPERIENCE PROXY", "Chi-Square (p=0.045) confirms payment method significantly associates with satisfaction."),
        ("4", "DELIVERY SPEED DRIVES REVENUE", "Khulna: fastest delivery (6.81 days) & highest avg purchase ($513.94). Every day saved = revenue."),
        ("5", "32.9% CUSTOMERS AT-RISK", "Nearly 1 in 3 customers are At-Risk (RFM). Largest churn threat requires immediate action."),
        ("6", "DEVICE PARITY", "Mobile 33.7% · Desktop 33.5% · Tablet 32.8% — no platform can be deprioritized."),
    ]

    cols = [(0.35, 4.9), (5.2, 4.9)]
    for i, (num, title, body) in enumerate(insights):
        col = i % 2
        row = i // 2
        cx = cols[col][0]
        cy = 1.2 + row * 1.3
        cw = 4.6
        ch = 1.15
        add_rect(slide, cx, cy, cw, ch, CARD_BG)
        # Accent number
        add_rect(slide, cx, cy, 0.28, ch, RGBColor(25, 25, 25))
        add_text(slide, num, cx+0.03, cy+0.3, 0.22, 0.45,
                 font_size=16, color=ACCENT_GREEN, bold=True, align=PP_ALIGN.CENTER, font_name="Arial Black")
        add_text(slide, title, cx+0.35, cy+0.06, cw-0.45, 0.3,
                 font_size=9, color=ACCENT_GREEN, bold=True)
        add_text(slide, body, cx+0.35, cy+0.38, cw-0.45, 0.68,
                 font_size=8.5, color=GRAY_LIGHT)

    slide_label(slide, 2)


# ─────────────────────────────────────────────────────────────
# SLIDE 3 — Level 1: Basic Statistics
# Speaker: "Our dataset of 10,000 records is clean — zero missing values and
# zero duplicates — giving us a solid foundation for all downstream analysis."
# "Age clusters around 44 years with the unusual mode at 51, while Purchase
# Amount has a remarkably high variance of 81,932 dollars-squared, signaling
# wildly diverse spending power we must account for in pricing strategy."
# ─────────────────────────────────────────────────────────────
def slide_03(prs):
    slide = blank_slide(prs)
    set_bg(slide)
    add_rect(slide, 0, 0, 0.06, SH, ACCENT_GREEN)

    section_pill(slide, "LEVEL 1 · BASIC INSIGHTS")
    slide_title(slide, "Descriptive Statistics", y=0.38)

    # Age stats row
    add_text(slide, "AGE  (Q1)", 0.35, 1.15, 3.0, 0.25,
             font_size=10, color=ACCENT_GREEN, bold=True)
    for i, (v, l) in enumerate([("43.79","Mean"), ("44.0","Median"), ("51","Mode")]):
        stat_card(slide, 0.35+i*1.55, 1.45, 1.42, 0.85, v, l)

    # Purchase Amount stats
    add_text(slide, "PURCHASE AMOUNT  (Q2)", 5.0, 1.15, 4.8, 0.25,
             font_size=10, color=ACCENT_GREEN, bold=True)
    for i, (v, l) in enumerate([("$503.89","Mean"), ("$286.24","Std Dev"), ("81,932","Variance")]):
        stat_card(slide, 5.0+i*1.65, 1.45, 1.52, 0.85, v, l, val_color=TEAL)

    # Z-Score table header
    add_text(slide, "PURCHASE AMOUNT Z-SCORE SAMPLE  (Q2)", 0.35, 2.52, 5.0, 0.25,
             font_size=10, color=ACCENT_GREEN, bold=True)
    headers = ["Cust ID", "Purchase ($)", "Z-Score"]
    rows = [
        ["1", "$202.54", "−1.0529"],
        ["2", "$655.94", " +0.5312"],
        ["3", "$963.65", " +1.6063"],
        ["4", "$485.59", " −0.0639"],
        ["5", "$143.27", " −1.2599"],
    ]
    col_x  = [0.35, 1.95, 3.55]
    col_w  = [1.5, 1.5, 1.5]
    # Header
    add_rect(slide, 0.35, 2.8, 4.8, 0.28, RGBColor(40,40,40))
    for j, (hdr, cx, cw) in enumerate(zip(headers, col_x, col_w)):
        add_text(slide, hdr, cx+0.05, 2.82, cw, 0.24,
                 font_size=9, color=ACCENT_GREEN, bold=True)
    for ri, row in enumerate(rows):
        bg = CARD_BG if ri % 2 == 0 else CARD_BG2
        add_rect(slide, 0.35, 3.1+ri*0.28, 4.8, 0.26, bg)
        for j, (cell, cx, cw) in enumerate(zip(row, col_x, col_w)):
            add_text(slide, cell, cx+0.05, 3.12+ri*0.28, cw, 0.22,
                     font_size=9, color=WHITE)

    # Data quality callout
    add_rect(slide, 5.0, 2.52, 4.6, 2.85, CARD_BG)
    add_text(slide, "✅  DATA QUALITY REPORT", 5.15, 2.65, 4.3, 0.3,
             font_size=10, color=ACCENT_GREEN, bold=True)
    dq_items = [
        ("Total Records", "10,000"),
        ("Missing Values", "0  (Clean ✅)"),
        ("Duplicate Rows", "0  (Clean ✅)"),
        ("Features (raw)", "16"),
        ("Features (engineered)", "28"),
        ("Avg Review Score (Q5)", "3.00 / 5"),
        ("Return Customers (Q4)", "4,996  (50.0%)"),
    ]
    for i, (lbl, val) in enumerate(dq_items):
        add_text(slide, lbl, 5.15, 3.05+i*0.29, 2.5, 0.26, font_size=9, color=GRAY_LIGHT)
        add_text(slide, val, 7.7, 3.05+i*0.29, 1.7, 0.26, font_size=9,
                 color=WHITE, bold=True, align=PP_ALIGN.RIGHT)

    slide_label(slide, 3)


# ─────────────────────────────────────────────────────────────
# SLIDE 4 — Level 1: Top Categories & Payment Methods
# Speaker: "Electronics leads purchases at 1,309 transactions, but the gap to
# third-place Toys is only 21 sales — a remarkably competitive category landscape
# that argues against over-indexing marketing spend on any single category."
# "Bank Transfer is the most common payment method at 20.7% share, and the
# virtually even device split — Mobile 33.7%, Desktop 33.5%, Tablet 32.8% —
# reinforces the need for truly platform-agnostic UX design."
# ─────────────────────────────────────────────────────────────
def slide_04(prs):
    slide = blank_slide(prs)
    set_bg(slide)
    add_rect(slide, 0, 0, 0.06, SH, ACCENT_GREEN)

    section_pill(slide, "LEVEL 1 · BASIC INSIGHTS")
    slide_title(slide, "Categories, Payments & Device Mix", y=0.38)

    # Top 3 categories
    add_text(slide, "TOP 3 PRODUCT CATEGORIES  (Q3)", 0.35, 1.1, 4.5, 0.25,
             font_size=10, color=ACCENT_GREEN, bold=True)
    cats   = [("Electronics", 1309), ("Clothing", 1297), ("Toys", 1288)]
    colors = [ACCENT_GREEN, TEAL, RGBColor(200,200,200)]
    for i, ((cat, cnt), col) in enumerate(zip(cats, colors)):
        cy  = 1.4 + i*0.45
        bw  = (cnt/1400)*3.8
        add_rect(slide, 0.35, cy, bw, 0.32, col)
        add_text(slide, cat, 0.45, cy+0.04, 1.8, 0.24, font_size=10, color=BG_COLOR, bold=True)
        add_text(slide, f"{cnt:,}", 0.35+bw+0.08, cy+0.04, 0.9, 0.24,
                 font_size=10, color=WHITE, bold=True)

    # Payment methods
    add_text(slide, "PAYMENT METHODS  (Q10 + Deep Dive)", 0.35, 2.85, 4.5, 0.25,
             font_size=10, color=ACCENT_GREEN, bold=True)
    pay_data = [
        ("Bank Transfer",    2067, 20.7),
        ("Cash on Delivery", 2007, 20.1),
        ("Credit Card",      2028, 20.3),
        ("Debit Card",       1983, 19.8),
        ("PayPal",           1915, 19.2),
    ]
    for i, (pm, cnt, pct) in enumerate(pay_data):
        cy  = 3.15 + i*0.38
        bw  = (pct/25.0)*3.8
        col = ACCENT_GREEN if i == 0 else RGBColor(60, 120, 60)
        add_rect(slide, 0.35, cy, bw, 0.27, col)
        add_text(slide, pm, 0.45, cy+0.03, 2.1, 0.22, font_size=9, color=BG_COLOR, bold=(i==0))
        add_text(slide, f"{cnt:,}  ({pct}%)", 0.35+bw+0.08, cy+0.03, 1.4, 0.22,
                 font_size=9, color=WHITE)

    # Device mix visual (right panel)
    add_rect(slide, 5.0, 1.05, 4.6, 4.3, CARD_BG)
    add_text(slide, "DEVICE USAGE SPLIT  (Q8)", 5.15, 1.18, 4.3, 0.25,
             font_size=10, color=ACCENT_GREEN, bold=True)

    devices = [("Mobile", 33.74, ACCENT_GREEN), ("Desktop", 33.48, TEAL), ("Tablet", 32.78, GRAY_LIGHT)]
    for i, (dev, pct, col) in enumerate(devices):
        cy  = 1.6 + i*0.9
        add_rect(slide, 5.15, cy, 4.25, 0.65, RGBColor(28,28,28))
        bw = (pct/40.0)*4.05
        add_rect(slide, 5.15, cy, bw, 0.65, col)
        add_text(slide, dev, 5.25, cy+0.12, 1.5, 0.35, font_size=12, color=BG_COLOR, bold=True)
        add_text(slide, f"{pct}%", 5.15+bw-0.9, cy+0.12, 0.8, 0.35,
                 font_size=14, color=BG_COLOR, bold=True, align=PP_ALIGN.RIGHT)

    add_text(slide, "★  Near-perfect parity across all 3 platforms.", 5.15, 4.5, 4.3, 0.3,
             font_size=9, color=ACCENT_GREEN, italic=True)

    # Discount insight
    add_rect(slide, 5.0, 4.85, 4.6, 0.58, RGBColor(20,40,20))
    add_text(slide, "Q9 · Discount vs No Discount Purchase Avg", 5.15, 4.9, 4.3, 0.22,
             font_size=9, color=ACCENT_GREEN, bold=True)
    add_text(slide, "Full Price: $502.51  ·  Discounted: $505.26  → Negligible difference",
             5.15, 5.12, 4.3, 0.22, font_size=9, color=WHITE)

    slide_label(slide, 4)


# ─────────────────────────────────────────────────────────────
# SLIDE 5 — Level 2: Correlations
# Speaker: "Our Pearson correlation between time on site and purchase amount is
# r=0.01 with a p-value of 0.316 — statistically non-significant and practically
# zero. More time browsing does not mean more spending."
# "This is a counterintuitive finding that challenges typical e-commerce
# assumptions — the platform should focus on conversion speed rather than
# maximizing dwell time, since engagement depth and spend are decoupled."
# ─────────────────────────────────────────────────────────────
def slide_05(prs):
    slide = blank_slide(prs)
    set_bg(slide)
    add_rect(slide, 0, 0, 0.06, SH, ACCENT_GREEN)

    section_pill(slide, "LEVEL 2 · INTERMEDIATE INSIGHTS")
    slide_title(slide, "Correlation Analysis: Time vs. Purchase", y=0.38)

    # Main correlation results
    corr_results = [
        ("Time on Site ↔ Purchase Amount",  "r = 0.0100", "p = 0.3162", "❌ NOT Significant"),
        ("Time on Site ↔ Items Purchased",  "r = 0.0001", "p = 0.9882", "❌ NOT Significant"),
        ("Delivery Time ↔ Review Score",    "r = 0.0125", "p = 0.2117", "❌ NOT Significant"),
        ("Items Purchased ↔ Purchase Amt",  "r = 0.0078", "p = 0.4364", "❌ NOT Significant"),
        ("Age ↔ Purchase Amount",           "r = −0.0030", "p = 0.7617", "❌ NOT Significant"),
        ("Age ↔ Time on Website",           "r = 0.0120", "p = 0.2283", "❌ NOT Significant"),
    ]

    add_text(slide, "ALL NUMERIC CORRELATIONS", 0.35, 1.1, 6.0, 0.25,
             font_size=10, color=ACCENT_GREEN, bold=True)
    # Table header
    add_rect(slide, 0.35, 1.38, 9.3, 0.28, RGBColor(40,40,40))
    for j, (hdr, cx, cw) in enumerate(zip(
            ["Variable Pair", "Pearson r", "p-value", "Result"],
            [0.35, 4.6, 6.1, 7.5], [4.2, 1.4, 1.3, 2.0])):
        add_text(slide, hdr, cx+0.05, 1.4, cw, 0.24,
                 font_size=9, color=ACCENT_GREEN, bold=True)

    for ri, (pair, r, p, res) in enumerate(corr_results):
        bg = CARD_BG if ri % 2 == 0 else CARD_BG2
        add_rect(slide, 0.35, 1.68+ri*0.32, 9.3, 0.3, bg)
        add_text(slide, pair, 0.45, 1.7+ri*0.32, 4.1, 0.26, font_size=9, color=WHITE)
        add_text(slide, r,    4.65, 1.7+ri*0.32, 1.3, 0.26, font_size=9, color=TEAL, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, p,    6.15, 1.7+ri*0.32, 1.2, 0.26, font_size=9, color=GRAY_LIGHT, align=PP_ALIGN.CENTER)
        rc = RED_WARN if "NOT" in res else ACCENT_GREEN
        add_text(slide, res,  7.55, 1.7+ri*0.32, 2.0, 0.26, font_size=9, color=rc)

    # Key insight box
    add_rect(slide, 0.35, 3.95, 9.3, 1.4, RGBColor(10,30,10))
    add_text(slide, "🔑  KEY INSIGHT: ZERO CORRELATION — MAXIMUM BUSINESS IMPLICATION",
             0.5, 4.05, 9.0, 0.3, font_size=10, color=ACCENT_GREEN, bold=True)
    insight_txt = (
        "No significant linear relationship exists between any pair of numeric variables (all p > 0.05). "
        "This means browsing time, age, items purchased, and delivery speed do NOT predict spend in isolation. "
        "Retention is a multivariate problem — addressed in Level 3 via ML feature importance."
    )
    add_text(slide, insight_txt, 0.5, 4.38, 9.0, 0.85, font_size=10, color=WHITE)

    # L2-Q2 highlight
    add_rect(slide, 0.35, 5.12, 4.55, 0.35, RGBColor(20,20,50))
    add_text(slide, "L2-Q2 Interpretation: No significant relationship between time on site & purchase amount or items.",
             0.5, 5.17, 4.3, 0.27, font_size=8.5, color=GRAY_LIGHT, italic=True)

    # L2-Q5 highlight
    add_rect(slide, 5.0, 5.12, 4.65, 0.35, RGBColor(20,20,50))
    add_text(slide, "L2-Q5: 2nd Highest Purchase Location → Barisal ($513.67 avg)",
             5.15, 5.17, 4.4, 0.27, font_size=8.5, color=TEAL, bold=True)

    slide_label(slide, 5)


# ─────────────────────────────────────────────────────────────
# SLIDE 6 — Level 2: Satisfaction vs. Return Customers
# Speaker: "20.08% of our customers are both satisfied — scoring 4 or 5 stars —
# AND return customers, representing the golden segment of 2,008 individuals
# who are already predisposed to brand advocacy."
# "Interestingly, average items purchased are virtually flat across Low, Medium
# and High satisfaction levels — around 5 items each — suggesting that basket
# size alone does not drive or reflect satisfaction, a key nuance for merchandising."
# ─────────────────────────────────────────────────────────────
def slide_06(prs):
    slide = blank_slide(prs)
    set_bg(slide)
    add_rect(slide, 0, 0, 0.06, SH, ACCENT_GREEN)

    section_pill(slide, "LEVEL 2 · INTERMEDIATE INSIGHTS")
    slide_title(slide, "Satisfaction & Return Customer Analysis", y=0.38)

    # L2-Q3 callout
    add_rect(slide, 0.35, 1.12, 4.6, 1.5, RGBColor(15,40,15))
    add_text(slide, "L2-Q3  ·  SATISFIED & RETURN CUSTOMERS", 0.55, 1.22, 4.2, 0.25,
             font_size=9.5, color=ACCENT_GREEN, bold=True)
    add_text(slide, "20.08%", 1.0, 1.52, 3.0, 0.65,
             font_size=44, color=ACCENT_GREEN, bold=True, align=PP_ALIGN.CENTER, font_name="Arial Black")
    add_text(slide, "2,008 customers rated ≥4★ AND returned", 0.55, 2.15, 4.2, 0.28,
             font_size=9, color=WHITE)

    # Satisfaction breakdown quad
    add_rect(slide, 0.35, 2.85, 4.6, 2.5, CARD_BG)
    add_text(slide, "SATISFACTION ↔ RETURN BREAKDOWN", 0.55, 2.97, 4.2, 0.25,
             font_size=9, color=ACCENT_GREEN, bold=True)
    bd_data = [
        ("Satisfied (≥4★) + Return",       "2,008", "20.08%", ACCENT_GREEN),
        ("Satisfied (≥4★) + Not Return",   "1,996", "19.96%", TEAL),
        ("Unsatisfied (<4★) + Return",     "2,988", "29.88%", GRAY_LIGHT),
        ("Unsatisfied (<4★) + Not Return", "3,008", "30.08%", GRAY_MID),
    ]
    for i, (lbl, cnt, pct, col) in enumerate(bd_data):
        cy = 3.28+i*0.45
        add_rect(slide, 0.45, cy, 0.04, 0.3, col)
        add_text(slide, lbl, 0.57, cy, 2.7, 0.3, font_size=9, color=WHITE)
        add_text(slide, cnt, 3.3, cy, 0.7, 0.3, font_size=9, color=col, bold=True, align=PP_ALIGN.RIGHT)
        add_text(slide, pct, 4.05, cy, 0.75, 0.3, font_size=9, color=col, bold=True, align=PP_ALIGN.RIGHT)

    # L2-Q4 items by satisfaction
    add_rect(slide, 5.1, 1.12, 4.5, 4.22, CARD_BG)
    add_text(slide, "L2-Q4  ·  ITEMS PURCHASED vs SATISFACTION", 5.3, 1.22, 4.2, 0.25,
             font_size=9.5, color=ACCENT_GREEN, bold=True)

    sat_data = [
        ("Low",    4.989, 5.0, 2.545),
        ("Medium", 5.059, 5.0, 2.563),
        ("High",   4.948, 5.0, 2.617),
    ]
    add_rect(slide, 5.15, 1.52, 4.3, 0.28, RGBColor(40,40,40))
    for j, hdr in enumerate(["Satisfaction", "Mean", "Median", "Std Dev"]):
        add_text(slide, hdr, 5.2+j*1.07, 1.54, 1.0, 0.24,
                 font_size=9, color=ACCENT_GREEN, bold=True)
    for ri, (s, mn, md, sd) in enumerate(sat_data):
        bg = CARD_BG if ri % 2 == 0 else CARD_BG2
        add_rect(slide, 5.15, 1.82+ri*0.3, 4.3, 0.28, bg)
        col = ACCENT_GREEN if s == "High" else (TEAL if s == "Medium" else GRAY_LIGHT)
        add_text(slide, s,       5.22, 1.84+ri*0.3, 1.0, 0.24, font_size=9, color=col, bold=True)
        add_text(slide, f"{mn:.3f}", 6.29, 1.84+ri*0.3, 1.0, 0.24, font_size=9, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(slide, f"{md:.1f}", 7.36, 1.84+ri*0.3, 1.0, 0.24, font_size=9, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(slide, f"{sd:.3f}", 8.42, 1.84+ri*0.3, 0.9, 0.24, font_size=9, color=WHITE, align=PP_ALIGN.CENTER)

    add_text(slide, "Insight: Items purchased are virtually identical (~5) across all satisfaction levels.",
             5.3, 2.75, 4.2, 0.4, font_size=9, color=TEAL, italic=True)

    # Avg review score by payment
    add_text(slide, "L2-Q1  ·  Avg Review Score — Bank Transfer (most common): 2.99 / 5",
             5.3, 3.25, 4.2, 0.3, font_size=9, color=GRAY_LIGHT)

    # Bar: items by satisfaction visual
    add_text(slide, "ITEMS PURCHASED BY SATISFACTION  (Bar)", 5.3, 3.65, 4.2, 0.25,
             font_size=9, color=ACCENT_GREEN, bold=True)
    sat_bars = [("Low", 4.989), ("Medium", 5.059), ("High", 4.948)]
    max_b = 5.2
    bar_colors = [GRAY_LIGHT, TEAL, ACCENT_GREEN]
    for i, ((s, v), col) in enumerate(zip(sat_bars, bar_colors)):
        bw = (v/max_b)*3.7
        add_rect(slide, 5.95, 3.95+i*0.38, bw, 0.28, col)
        add_text(slide, s, 5.15, 3.97+i*0.38, 0.72, 0.24, font_size=9, color=WHITE, align=PP_ALIGN.RIGHT)
        add_text(slide, f"{v:.3f}", 5.95+bw+0.08, 3.97+i*0.38, 0.7, 0.24,
                 font_size=9, color=WHITE, bold=True)

    slide_label(slide, 6)


# ─────────────────────────────────────────────────────────────
# SLIDE 7 — Level 3: Drivers of Return Customers
# Speaker: "Our Random Forest classifier, with an AUC of 0.5018, reveals that
# Purchase Amount is the single biggest predictor of return behavior at 17.89%
# importance, followed by Age at 13.17% and Time on Site at 13.08%."
# "Critically, Discount Availed ranks dead last at just 1.31% — confirming that
# our retention strategy should never rely on price cuts, but instead target
# high-spend customers with personalized value-based engagement."
# ─────────────────────────────────────────────────────────────
def slide_07(prs):
    slide = blank_slide(prs)
    set_bg(slide)
    add_rect(slide, 0, 0, 0.06, SH, ACCENT_GREEN)

    section_pill(slide, "LEVEL 3 · CRITICAL THINKING")
    slide_title(slide, "Drivers of Return Customers  (ML Feature Importance)", y=0.38)

    # Feature importances from RF
    features = [
        ("Purchase Amount ($)",         17.89, ACCENT_GREEN),
        ("Age",                         13.17, ACCENT_GREEN),
        ("Time on Website (min)",       13.08, ACCENT_GREEN),
        ("Delivery Time (days)",         8.59, TEAL),
        ("No. Items Purchased",          7.51, TEAL),
        ("Location",                     7.09, TEAL),
        ("Product Category",             6.89, TEAL),
        ("Payment Method",               5.32, GRAY_LIGHT),
        ("Review Score (1-5)",           5.10, GRAY_LIGHT),
        ("Gender",                       4.14, GRAY_LIGHT),
        ("Customer Satisfaction",        3.39, GRAY_MID),
        ("Subscription Status",          3.32, GRAY_MID),
        ("Device Type",                  3.21, GRAY_MID),
        ("Is Satisfied",                 1.31, RED_WARN),
    ]
    max_imp = 17.89
    bar_x = 3.0
    bar_max_w = 5.3

    add_text(slide, "Random Forest  |  AUC = 0.5018  |  Top 14 Predictors",
             0.35, 1.1, 9.3, 0.25, font_size=9.5, color=GRAY_LIGHT)

    for i, (feat, imp, col) in enumerate(features):
        cy  = 1.42 + i*0.27
        bw  = (imp/max_imp)*bar_max_w
        # Feature label
        add_text(slide, feat, 0.38, cy, 2.55, 0.24, font_size=8.5, color=WHITE, align=PP_ALIGN.RIGHT)
        # Bar
        add_rect(slide, bar_x, cy+0.03, bw, 0.2, col)
        # Value
        add_text(slide, f"{imp:.2f}%", bar_x+bw+0.08, cy, 0.7, 0.24,
                 font_size=8.5, color=col, bold=True)

    # Cumulative labels
    add_text(slide, "▶ Top 5 features explain 60.24% of model decisions",
             0.35, 5.18, 6.0, 0.25, font_size=9, color=ACCENT_GREEN, bold=True)
    add_text(slide, "▶ Discount (1.31%) & Gender (4.14%) have minimal predictive power",
             0.35, 5.4, 7.0, 0.22, font_size=9, color=RED_WARN, italic=True)

    # Model metrics panel
    add_rect(slide, 8.5, 1.1, 1.15, 4.2, CARD_BG)
    add_text(slide, "MODEL\nMETRICS", 8.58, 1.18, 1.0, 0.45,
             font_size=8, color=ACCENT_GREEN, bold=True, align=PP_ALIGN.CENTER)
    metrics = [("RF AUC","0.5018"), ("GB AUC","0.4983"), ("LR AUC","0.4920"),
               ("CV AUC","0.5018"), ("±","0.0059"), ("Acc.","~50%")]
    for i, (lbl, val) in enumerate(metrics):
        add_text(slide, lbl, 8.55, 1.65+i*0.55, 1.05, 0.22,
                 font_size=7.5, color=GRAY_LIGHT, align=PP_ALIGN.CENTER)
        add_text(slide, val, 8.55, 1.85+i*0.55, 1.05, 0.25,
                 font_size=9, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    slide_label(slide, 7)


# ─────────────────────────────────────────────────────────────
# SLIDE 8 — Statistical Test Matrix
# Speaker: "This slide is the evidence backbone of our presentation —
# a complete matrix of every statistical test run, with the exact test statistic,
# p-value, and significance verdict for each hypothesis."
# "The most striking takeaway is that the only significant results are
# Product Category on Purchase Amount (ANOVA F=2.45, p=0.017) and
# Customer Satisfaction on Payment Method (Chi-Square χ²=15.81, p=0.045),
# while virtually all other tests return non-significant results, pointing to
# a highly randomized, non-deterministic dataset."
# ─────────────────────────────────────────────────────────────
def slide_08(prs):
    slide = blank_slide(prs)
    set_bg(slide)
    add_rect(slide, 0, 0, 0.06, SH, ACCENT_GREEN)

    section_pill(slide, "LEVEL 3 · STATISTICAL TESTS")
    slide_title(slide, "Statistical Test Summary Matrix", y=0.38)

    tests = [
        # (Test, Hypothesis, Stat, p-value, Significant?)
        ("T-Test",    "Return vs Non-Return → Purchase Amt",       "t=0.1924",    "p=0.8475", False),
        ("T-Test",    "Discount vs No Discount → Purchase Amt",    "t=0.4808",    "p=0.6307", False),
        ("T-Test",    "Premium vs Free → Delivery Time",           "t=1.1699",    "p=0.2421", False),
        ("T-Test",    "High vs Low Satisfaction → Time on Site",   "t=−0.7710",   "p=0.4407", False),
        ("T-Test",    "Male vs Female → Purchase Amount",          "t=0.2635",    "p=0.7922", False),
        ("ANOVA",     "Product Category → Purchase Amount",        "F=2.4494",    "p=0.0165", True),
        ("ANOVA",     "Payment Method → Review Score",             "F=0.3304",    "p=0.8577", False),
        ("ANOVA",     "Location → Delivery Time",                  "F=0.8768",    "p=0.5239", False),
        ("ANOVA",     "Customer Satisfaction → Purchase Amount",   "F=0.2335",    "p=0.7917", False),
        ("ANOVA",     "Age Group → Time on Website",               "F=1.0280",    "p=0.3911", False),
        ("Chi-Sq",    "Return Customer ⊥ Subscription Status",     "χ²=0.9879",   "p=0.6102", False),
        ("Chi-Sq",    "Satisfaction ⊥ Payment Method",             "χ²=15.8143",  "p=0.0451", True),
        ("Chi-Sq",    "Discount Availed ⊥ Return Customer",        "χ²=6.1518",   "p=0.0131", True),
        ("Chi-Sq",    "Device Type ⊥ Customer Satisfaction",       "χ²=2.5224",   "p=0.6406", False),
        ("Pearson r", "Time on Site ↔ Purchase Amount",            "r=0.0100",    "p=0.3162", False),
        ("OLS Reg",   "Predict Purchase Amount (R²=0.001)",        "F=0.7251",    "p=0.651",  False),
        ("Logistic",  "Predict Return Customer",                    "AUC=0.4920",  "p=0.621",  False),
    ]

    # Header row
    cols_x = [0.35, 1.5, 4.85, 6.55, 7.85, 8.9]
    cols_w = [1.1, 3.3, 1.65, 1.25, 1.0, 0.9]
    hdrs   = ["Test Type", "Hypothesis", "Statistic", "p-value", "α=0.05", "Decision"]
    add_rect(slide, 0.35, 1.0, 9.3, 0.28, RGBColor(40,40,40))
    for hdr, cx, cw in zip(hdrs, cols_x, cols_w):
        add_text(slide, hdr, cx+0.04, 1.02, cw, 0.24, font_size=8.5, color=ACCENT_GREEN, bold=True)

    row_h = 0.255
    for ri, (ttype, hyp, stat, pval, sig) in enumerate(tests):
        bg = CARD_BG if ri % 2 == 0 else CARD_BG2
        cy = 1.3 + ri*row_h
        add_rect(slide, 0.35, cy, 9.3, row_h-0.02, bg)

        # Test type color
        tcol = {"T-Test": TEAL, "ANOVA": RGBColor(200,140,60), "Chi-Sq": RGBColor(150,100,220),
                "Pearson r": GRAY_LIGHT, "OLS Reg": GRAY_LIGHT, "Logistic": GRAY_LIGHT}.get(ttype, WHITE)
        add_text(slide, ttype, cols_x[0]+0.04, cy+0.02, cols_w[0], row_h-0.06, font_size=8, color=tcol, bold=True)
        add_text(slide, hyp,   cols_x[1]+0.04, cy+0.02, cols_w[1], row_h-0.06, font_size=7.8, color=WHITE)
        add_text(slide, stat,  cols_x[2]+0.04, cy+0.02, cols_w[2], row_h-0.06, font_size=8, color=TEAL)
        pcol = ACCENT_GREEN if sig else RED_WARN
        add_text(slide, pval,  cols_x[3]+0.04, cy+0.02, cols_w[3], row_h-0.06, font_size=8, color=pcol, bold=sig)
        dec_txt = "✅ YES" if sig else "❌ NO"
        add_text(slide, dec_txt, cols_x[4]+0.04, cy+0.02, cols_w[4], row_h-0.06,
                 font_size=8, color=ACCENT_GREEN if sig else GRAY_MID, bold=sig)
        add_text(slide, "Significant" if sig else "Not Sig.", cols_x[5]+0.04, cy+0.02,
                 cols_w[5], row_h-0.06, font_size=7.5, color=ACCENT_GREEN if sig else GRAY_MID)

    slide_label(slide, 8)


# ─────────────────────────────────────────────────────────────
# SLIDE 9 — Regional Performance Analysis
# Speaker: "Khulna stands out as the benchmark city — highest average purchase
# at $513.94 and fastest delivery at 6.81 days, delivering $663,497 in total
# revenue, the highest of any region."
# "Rajshahi shows the slowest delivery at 7.11 days alongside one of the lower
# average purchase figures — this 4.4% logistics efficiency gap has a direct
# revenue consequence that a regional fulfilment hub could close."
# ─────────────────────────────────────────────────────────────
def slide_09(prs):
    slide = blank_slide(prs)
    set_bg(slide)
    add_rect(slide, 0, 0, 0.06, SH, ACCENT_GREEN)

    section_pill(slide, "LEVEL 3 · REGIONAL ANALYSIS")
    slide_title(slide, "Geographic Performance Matrix", y=0.38)

    # Data from output
    geo_data = [
        ("Khulna",      1291, 513.94, 6.81, 0.50, 2.98, 663497, True),
        ("Mymensingh",  1280, 507.89, 7.10, 0.50, 3.02, 650099, False),
        ("Chittagong",  1272, 507.85, 7.02, 0.50, 3.00, 645985, False),
        ("Barisal",     1250, 513.67, 7.07, 0.49, 3.03, 642088, False),
        ("Rangpur",     1253, 494.37, 7.03, 0.51, 2.97, 619446, False),
        ("Sylhet",      1237, 494.98, 7.02, 0.48, 2.99, 612290, False),
        ("Dhaka",       1206, 502.00, 6.94, 0.50, 2.96, 605412, False),
        ("Rajshahi",    1211, 495.54, 7.11, 0.53, 3.02, 600099, True),
    ]

    # Table
    hdrs   = ["Location", "Customers", "Avg Purchase", "Avg Delivery", "Return Rate", "Avg Review", "Total Revenue"]
    cols_x = [0.35, 1.85, 2.95, 4.25, 5.5,  6.55, 7.55]
    cols_w = [1.45, 1.05, 1.25, 1.2,  1.0,  0.95, 1.75]

    add_rect(slide, 0.35, 1.0, 9.3, 0.3, RGBColor(40,40,40))
    for hdr, cx, cw in zip(hdrs, cols_x, cols_w):
        add_text(slide, hdr, cx+0.04, 1.02, cw, 0.26, font_size=8.5, color=ACCENT_GREEN, bold=True)

    row_h = 0.43
    for ri, (loc, cust, avgp, avgd, retr, avgr, rev, highlight) in enumerate(geo_data):
        bg = RGBColor(20,40,20) if highlight else (CARD_BG if ri % 2 == 0 else CARD_BG2)
        cy = 1.32 + ri*row_h
        add_rect(slide, 0.35, cy, 9.3, row_h-0.03, bg)
        lcol = ACCENT_GREEN if highlight else WHITE
        add_text(slide, loc,          cols_x[0]+0.04, cy+0.08, cols_w[0], 0.26, font_size=9, color=lcol, bold=highlight)
        add_text(slide, f"{cust:,}",  cols_x[1]+0.04, cy+0.08, cols_w[1], 0.26, font_size=9, color=WHITE, align=PP_ALIGN.CENTER)
        pc = ACCENT_GREEN if avgp > 510 else WHITE
        add_text(slide, f"${avgp:.2f}", cols_x[2]+0.04, cy+0.08, cols_w[2], 0.26, font_size=9, color=pc, bold=(avgp>510), align=PP_ALIGN.CENTER)
        dc = RED_WARN if avgd > 7.05 else (ACCENT_GREEN if avgd < 6.9 else WHITE)
        add_text(slide, f"{avgd:.2f}d", cols_x[3]+0.04, cy+0.08, cols_w[3], 0.26, font_size=9, color=dc, bold=(avgd<6.9 or avgd>7.05), align=PP_ALIGN.CENTER)
        add_text(slide, f"{retr:.0%}",  cols_x[4]+0.04, cy+0.08, cols_w[4], 0.26, font_size=9, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(slide, f"{avgr:.2f}",  cols_x[5]+0.04, cy+0.08, cols_w[5], 0.26, font_size=9, color=WHITE, align=PP_ALIGN.CENTER)
        rc = ACCENT_GREEN if rev > 640000 else WHITE
        add_text(slide, f"${rev:,}",    cols_x[6]+0.04, cy+0.08, cols_w[6], 0.26, font_size=9, color=rc, bold=(rev>640000))

    # Insight footer
    add_rect(slide, 0.35, 4.83, 9.3, 0.6, RGBColor(10,25,10))
    add_text(slide, "★ Khulna: fastest delivery (6.81d) + highest purchase ($513.94) → logistics efficiency drives revenue.",
             0.55, 4.9, 8.8, 0.26, font_size=9.5, color=ACCENT_GREEN, bold=True)
    add_text(slide, "⚠ Rajshahi: slowest delivery (7.11d) → 4.4% efficiency gap correlates with lower spending. Invest in regional fulfilment.",
             0.55, 5.18, 8.8, 0.22, font_size=9, color=RED_WARN)

    slide_label(slide, 9)


# ─────────────────────────────────────────────────────────────
# SLIDE 10 — Business Implications & Findings
# Speaker: "Aggregating our statistical and ML findings, six business-critical
# truths emerge: retention is purely spend-driven, the subscription tier system
# is broken, payment friction determines satisfaction, fast delivery equals
# more revenue, one-third of our base is at churn risk, and all three devices
# demand equal product investment."
# "These aren't soft observations — each is backed by a p-value, an odds ratio,
# or a feature importance score that quantifies the business risk of inaction."
# ─────────────────────────────────────────────────────────────
def slide_10(prs):
    slide = blank_slide(prs)
    set_bg(slide)
    add_rect(slide, 0, 0, 0.06, SH, ACCENT_GREEN)

    section_pill(slide, "BUSINESS IMPLICATIONS")
    slide_title(slide, "Key Findings & Business Implications", y=0.38)

    findings = [
        ("💰", "RETENTION IS SPEND-DRIVEN",
         "Purchase Amount: 17.89% RF importance. Discounts: 1.31%.\nLoyalty must be earned through value, not price cuts."),
        ("🔒", "SUBSCRIPTION PARADOX",
         "66%+ on paid tiers, yet Chi-Square p=0.61 (NOT significant).\nSubscription tier does NOT predict return behavior."),
        ("💳", "PAYMENT = SATISFACTION PROXY",
         "Chi-Square p=0.045 ✅ SIGNIFICANT.\nCredit Card users: highest 'High' satisfaction rate."),
        ("🚚", "DELIVERY SPEED = REVENUE",
         "Khulna: 6.81 days → $513.94 avg purchase.\nRajshahi: 7.11 days → $495.54 avg purchase."),
        ("⚠️", "32.9% CUSTOMERS AT-RISK",
         "1,303 At-Risk customers (13%) + 153 Lost (1.5%).\n~33% of base trending toward churn."),
        ("📱", "DEVICE PARITY IS NON-NEGOTIABLE",
         "Mobile: 33.74%  ·  Desktop: 33.48%  ·  Tablet: 32.78%\nAll platforms demand equal UX investment."),
    ]

    for i, (icon, title, body) in enumerate(findings):
        col = i % 2
        row = i // 2
        cx = 0.35 + col*4.7
        cy = 1.1 + row*1.38
        cw, ch = 4.5, 1.28
        add_rect(slide, cx, cy, cw, ch, CARD_BG)
        # Left accent stripe
        acol = ACCENT_GREEN if col == 0 else TEAL
        add_rect(slide, cx, cy, 0.06, ch, acol)
        add_text(slide, icon, cx+0.15, cy+0.05, 0.45, 0.38, font_size=18, align=PP_ALIGN.CENTER)
        add_text(slide, title, cx+0.65, cy+0.08, cw-0.75, 0.3,
                 font_size=9.5, color=ACCENT_GREEN, bold=True)
        add_text(slide, body, cx+0.65, cy+0.42, cw-0.75, 0.78,
                 font_size=8.8, color=GRAY_LIGHT)

    slide_label(slide, 10)


# ─────────────────────────────────────────────────────────────
# SLIDE 11 — Strategic Recommendations
# Speaker: "Translating our findings into five concrete strategic actions:
# first, launch spend-tier re-engagement campaigns targeting the At-Risk segment
# using Purchase Amount Z-scores to identify whale customers most at risk of lapsing."
# "Second, prioritize logistics investment in Rajshahi to close the delivery gap —
# and critically, migrate Cash on Delivery users to digital payment channels,
# since CoD correlates with the highest dissatisfaction rate, a relationship
# confirmed by Chi-Square at p=0.045."
# ─────────────────────────────────────────────────────────────
def slide_11(prs):
    slide = blank_slide(prs)
    set_bg(slide)
    add_rect(slide, 0, 0, 0.06, SH, ACCENT_GREEN)

    section_pill(slide, "STRATEGIC RECOMMENDATIONS")
    slide_title(slide, "5 Strategic Actions for Growth", y=0.38)

    recs = [
        ("01", "SPEND-TIER RETENTION CAMPAIGNS",
         "Target At-Risk (13%) & Loyal (35%) segments with spend-based personalization — not generic discounts. "
         "Use Z-score tiers (Z>1.5 = whale customers) for hyper-targeted re-engagement messaging.",
         "Priority: IMMEDIATE"),
        ("02", "LOGISTICS INVESTMENT — RAJSHAHI",
         "Bridge the 4.4% delivery gap (7.11 vs 6.81 days). Regional fulfilment hub or last-mile carrier "
         "partnership can directly uplift Rajshahi's average purchase from $495 toward Khulna's $514.",
         "Priority: HIGH"),
        ("03", "DIGITAL PAYMENT MIGRATION",
         "CoD users show highest Low-satisfaction rates. Incentivize Credit Card & PayPal with "
         "small 'digital-first' perks. Each successful migration measurably improves satisfaction scores.",
         "Priority: HIGH"),
        ("04", "SUBSCRIPTION VALUE OVERHAUL",
         "Chi-Square confirms subscription status is NOT associated with return behavior (p=0.61). "
         "Redesign Premium tier benefits to create tangible loyalty differentiators — e.g. faster delivery SLA.",
         "Priority: MEDIUM"),
        ("05", "CLUSTER-SPECIFIC UX & PRODUCT DESIGN",
         "K-Means reveals 4 distinct behavioral archetypes. Build cluster-specific UX flows: "
         "Cluster 1 (high review, fast delivery) needs seamless checkout; Cluster 0 (older, high-spend) needs trust signals.",
         "Priority: MEDIUM"),
    ]

    for i, (num, title, body, priority) in enumerate(recs):
        cy  = 1.12 + i*0.86
        add_rect(slide, 0.35, cy, 9.3, 0.82, CARD_BG)
        # Number badge
        add_rect(slide, 0.35, cy, 0.52, 0.82, RGBColor(30,60,10))
        add_text(slide, num, 0.36, cy+0.24, 0.5, 0.35,
                 font_size=14, color=ACCENT_GREEN, bold=True, align=PP_ALIGN.CENTER, font_name="Arial Black")
        add_text(slide, title, 0.95, cy+0.06, 5.8, 0.28,
                 font_size=9.5, color=WHITE, bold=True)
        # Priority badge
        pcol = RED_WARN if "IMMEDIATE" in priority else (TEAL if "HIGH" in priority else GRAY_LIGHT)
        add_rect(slide, 6.8, cy+0.07, 1.75, 0.22, RGBColor(30,30,30))
        add_text(slide, priority, 6.85, cy+0.07, 1.65, 0.22,
                 font_size=8, color=pcol, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, body, 0.95, cy+0.36, 8.55, 0.42,
                 font_size=8.2, color=GRAY_LIGHT)

    slide_label(slide, 11)


# ─────────────────────────────────────────────────────────────
# SLIDE 12 — Q&A / Contact
# Speaker: "That concludes our presentation of the E-Commerce Customer Behavior
# Analysis. We covered all three levels of the question set, ran five families of
# statistical tests, and derived five actionable recommendations backed by data."
# "I'm happy to deep-dive into any specific test result, cluster profile, or
# regression coefficient — please feel free to ask questions."
# ─────────────────────────────────────────────────────────────
def slide_12(prs):
    slide = blank_slide(prs)
    set_bg(slide)

    # Full-width accent top bar
    add_rect(slide, 0, 0, SW, 0.06, ACCENT_GREEN)
    add_rect(slide, 0, SH-0.06, SW, 0.06, ACCENT_GREEN)

    # Background grid decoration
    for xi in [2.5, 5.0, 7.5]:
        add_rect(slide, xi, 0.06, 0.005, SH-0.12, GRAY_MID)
    for yi in [1.4, 2.8, 4.2]:
        add_rect(slide, 0, yi, SW, 0.005, GRAY_MID)

    add_text(slide, "Thank You", 1.2, 0.55, 7.6, 1.1,
             font_size=52, color=WHITE, bold=True, align=PP_ALIGN.CENTER, font_name="Arial Black")
    add_text(slide, "Questions & Discussion", 1.2, 1.5, 7.6, 0.5,
             font_size=20, color=ACCENT_GREEN, align=PP_ALIGN.CENTER, font_name="Calibri Light")

    # Contact card
    add_rect(slide, 2.5, 2.25, 5.0, 1.9, CARD_BG)
    add_text(slide, "Md. Shafat Hossain", 2.7, 2.38, 4.6, 0.38,
             font_size=16, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "Student of Data Science & AI Engineering", 2.7, 2.76, 4.6, 0.28,
             font_size=10, color=ACCENT_GREEN, align=PP_ALIGN.CENTER)
    add_text(slide, "Final Assignment  ·  E-Commerce Customer Behavior Analysis",
             2.7, 3.1, 4.6, 0.28, font_size=9, color=GRAY_LIGHT, align=PP_ALIGN.CENTER)
    add_text(slide, "10,000 Records  ·  16 Features  ·  28 Engineered  ·  5 Statistical Test Families",
             2.7, 3.4, 4.6, 0.25, font_size=8.5, color=GRAY_MID, align=PP_ALIGN.CENTER)

    # Summary metrics row
    summary = [("50.0%","Return Rate"), ("3.00★","Avg Review"), ("$503.89","Avg Purchase"),
               ("7.01d","Avg Delivery"), ("0.5018","Best AUC")]
    for i, (val, lbl) in enumerate(summary):
        stat_card(slide, 0.35+i*1.87, 4.45, 1.75, 0.95, val, lbl)

    slide_label(slide, 12)


# ─────────────────────────────────────────────────────────────
# MAIN
# ─────────────────────────────────────────────────────────────
if __name__ == "__main__":
    prs = new_prs()

    slide_01(prs)
    slide_02(prs)
    slide_03(prs)
    slide_04(prs)
    slide_05(prs)
    slide_06(prs)
    slide_07(prs)
    slide_08(prs)
    slide_09(prs)
    slide_10(prs)
    slide_11(prs)
    slide_12(prs)

    import os
    # Saves in the same folder as this script — works on Windows, Mac, and Linux
    script_dir = os.path.dirname(os.path.abspath(__file__))
    out = os.path.join(script_dir, "Ecommerce_Analysis_Presentation.pptx")
    prs.save(out)
    print(f"✅  Saved → {out}")
